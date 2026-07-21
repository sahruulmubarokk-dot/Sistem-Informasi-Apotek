import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

def get_fitted_dimensions(img_path, max_w_in, max_h_in):
    try:
        with Image.open(img_path) as img:
            orig_w, orig_h = img.size
        aspect = orig_w / orig_h
        if (max_w_in / max_h_in) > aspect:
            fit_h = max_h_in
            fit_w = max_h_in * aspect
        else:
            fit_w = max_w_in
            fit_h = max_w_in / aspect
        return Inches(fit_w), Inches(fit_h)
    except Exception:
        return Inches(max_w_in), Inches(max_h_in)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Color Palette (Vibrant Cyber Teal & Deep Navy Dark Mode)
    COLOR_BG = RGBColor(11, 19, 36)          # #0B1324 - Deep Slate Navy
    COLOR_CARD_BG = RGBColor(21, 34, 56)     # #152238 - Dark Card Surface
    COLOR_CARD_BORDER = RGBColor(38, 58, 92) # #263A5C - Subtle Border
    COLOR_CYAN = RGBColor(14, 165, 233)      # #0EA5E9 - Vibrant Cyan Accent
    COLOR_TEAL = RGBColor(45, 212, 191)      # #2DD4BF - Mint Teal Accent
    COLOR_EMERALD = RGBColor(52, 211, 153)   # #34D399 - Cyber Emerald Accent
    COLOR_WHITE = RGBColor(248, 250, 252)    # #F8FAFC - Pure Soft White
    COLOR_MUTED = RGBColor(148, 163, 184)    # #94A3B8 - Muted Silver Gray
    COLOR_ACCENT_BG = RGBColor(30, 41, 67)   # #1E2943 - Highlight Card Fill
    COLOR_FOOTER = RGBColor(100, 116, 139)   # #64748B - Footer Text

    def set_bg(slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_BG
        bg_shape.line.fill.background()

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.2)
        else:
            shape.line.fill.background()
        return shape

    def add_header(slide, category_text, title_text, slide_num=None):
        # Category Tag / Pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = f"APOTEKSIM  |  {category_text.upper()}"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_TEAL

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.733), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        # Subtle Accent Header Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_CARD_BORDER
        line.line.fill.background()

        # Footer
        ft_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.733), Inches(0.35))
        tf_ft = ft_box.text_frame
        tf_ft.margin_left = tf_ft.margin_right = tf_ft.margin_top = tf_ft.margin_bottom = 0
        p_ft = tf_ft.paragraphs[0]
        p_ft.text = "Universitas Dr. Soetomo Surabaya • Teknik Informatika • Kelompok 4 (2026)"
        p_ft.font.size = Pt(9)
        p_ft.font.color.rgb = COLOR_FOOTER
        
        if slide_num:
            p_ft2 = tf_ft.add_paragraph()
            p_ft2.text = f"Slide {slide_num} / 8"
            p_ft2.font.size = Pt(9)
            p_ft2.font.color.rgb = COLOR_TEAL
            p_ft2.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE & TEAM
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)

    # Decorative background card container
    add_card(slide1, Inches(0.8), Inches(0.55), Inches(11.733), Inches(6.35), bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)

    # Top Category Tag
    pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(0.85), Inches(4.5), Inches(0.35))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLOR_ACCENT_BG
    pill.line.color.rgb = COLOR_TEAL
    pill.line.width = Pt(1)
    tf_pill = pill.text_frame
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "LAPORAN AKHIR PROJECT-BASED LEARNING (PBL)"
    p_pill.font.size = Pt(9.5)
    p_pill.font.bold = True
    p_pill.font.color.rgb = COLOR_TEAL
    p_pill.alignment = PP_ALIGN.CENTER

    # Project Title & Subtitle
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(8.0), Inches(1.3))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t1 = tf_t.paragraphs[0]
    p_t1.text = "ApotekSim"
    p_t1.font.size = Pt(42)
    p_t1.font.bold = True
    p_t1.font.color.rgb = COLOR_CYAN

    p_t2 = tf_t.add_paragraph()
    p_t2.text = "Sistem Informasi Apotek Berbasis Web Modern dengan Arsitektur Full-Stack & Integrasi Supabase"
    p_t2.font.size = Pt(15)
    p_t2.font.color.rgb = COLOR_WHITE
    p_t2.space_before = Pt(4)

    # Logo Unitomo (Proportional fit inside Card)
    if os.path.exists("UNITOMO.jpg"):
        add_card(slide1, Inches(9.7), Inches(0.85), Inches(2.4), Inches(2.25), bg_color=RGBColor(255, 255, 255), border_color=COLOR_TEAL)
        fw, fh = get_fitted_dimensions("UNITOMO.jpg", 2.1, 1.95)
        # Center inside logo card
        lx = Inches(9.7) + (Inches(2.4) - fw) / 2
        ly = Inches(0.85) + (Inches(2.25) - fh) / 2
        slide1.shapes.add_picture("UNITOMO.jpg", lx, ly, fw, fh)

    # University Header
    univ_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.75), Inches(10.5), Inches(0.4))
    tf_u = univ_box.text_frame
    p_u = tf_u.paragraphs[0]
    p_u.text = "PROGRAM STUDI TEKNIK INFORMATIKA  •  FAKULTAS TEKNIK  •  UNIVERSITAS DR. SOETOMO SURABAYA"
    p_u.font.size = Pt(10.5)
    p_u.font.bold = True
    p_u.font.color.rgb = COLOR_MUTED

    # Divider Line
    div = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.2), Inches(10.933), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_CARD_BORDER
    div.line.fill.background()

    # Team Members Section Label
    team_lbl = slide1.shapes.add_textbox(Inches(1.2), Inches(3.32), Inches(5.0), Inches(0.3))
    p_tl = team_lbl.text_frame.paragraphs[0]
    p_tl.text = "TIM PENGEMBANG (KELOMPOK 4):"
    p_tl.font.size = Pt(11)
    p_tl.font.bold = True
    p_tl.font.color.rgb = COLOR_TEAL

    # Team Member Cards (3 Columns)
    team_data = [
        {"name": "Fahri Adis Al Hafni", "nim": "NIM: 202511420105", "role": "Fullstack Developer & DB Designer", "highlight": True},
        {"name": "Sahrul Mubarok", "nim": "NIM: 202511420107", "role": "Backend & Cloud Infrastructure", "highlight": False},
        {"name": "Ach Tohir", "nim": "NIM: 202511420048", "role": "Frontend UI/UX & QA Tester", "highlight": False}
    ]

    card_w = Inches(3.44)
    gap = Inches(0.3)
    start_x = Inches(1.2)

    for i, member in enumerate(team_data):
        cx = start_x + i * (card_w + gap)
        cy = Inches(3.75)
        ch = Inches(2.7)

        add_card(slide1, cx, cy, card_w, ch, bg_color=COLOR_ACCENT_BG if member["highlight"] else COLOR_BG, border_color=COLOR_CYAN if member["highlight"] else COLOR_CARD_BORDER)

        tb = slide1.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.25), card_w - Inches(0.4), ch - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = member["name"]
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = member["nim"]
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_CYAN
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = f"Peran Utama:\n{member['role']}"
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_MUTED
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 2: PROBLEM STATEMENT
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "ANALISIS OPERASIONAL & RISET MASALAH", "Problem Statement: Kelemahan Sistem Apotek Konvensional", 2)

    problems = [
        {
            "icon": "🚨",
            "title": "Kerugian Finansial (Obat Expired)",
            "desc": "Ketiadaan pencatatan berbasis batch & kedaluwarsa berakibat penumpukan obat ED di gudang tanpa terdeteksi sejak dini. Obat tidak dapat diretur ke supplier."
        },
        {
            "icon": "📉",
            "title": "Human Error & Selisih Stok",
            "desc": "Pencatatan manual berbasis buku / spreadsheet lokal tidak terintegrasi dengan transaksi kasir harian, memicu selisih stok fisik di gudang secara berulang."
        },
        {
            "icon": "⚖️",
            "title": "Risiko Kepatuhan Regulasi Medis",
            "desc": "Penjualan obat keras / resep dokter tanpa pengawasan data Dokter & No. SIP yang tervalidasi berpotensi sanksi administratif hingga pencabutan izin operasional."
        },
        {
            "icon": "🛡️",
            "title": "Kerentanan Arsitektur Legacy",
            "desc": "Sistem monolitik tradisional lambat, sulit dikembangkan, dan rentan terhadap serangan manipulasi data dari sisi klien (client-side manipulation)."
        }
    ]

    grid_w = Inches(5.65)
    grid_h = Inches(2.35)

    for i, p in enumerate(problems):
        row = i // 2
        col = i % 2
        gx = Inches(0.8) + col * (grid_w + Inches(0.433))
        gy = Inches(1.55) + row * (grid_h + Inches(0.3))

        add_card(slide2, gx, gy, grid_w, grid_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)

        tb = slide2.shapes.add_textbox(gx + Inches(0.25), gy + Inches(0.2), grid_w - Inches(0.5), grid_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"{p['icon']}  {p['title']}"
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CYAN

        p2 = tf.add_paragraph()
        p2.text = p['desc']
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(8)

    # Bottom Stat / Callout Banner
    add_card(slide2, Inches(0.8), Inches(6.35), Inches(11.733), Inches(0.6), bg_color=COLOR_ACCENT_BG, border_color=COLOR_TEAL)
    tb_bot = slide2.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5))
    tf_b = tb_bot.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "🎯 SOLUSI APOTEKSIM: Mengeliminasi Kerugian Kedaluwarsa, Menjamin Kecepatan Kasir & Keamanan Data Cloud"
    p_b.font.size = Pt(11.5)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_EMERALD
    p_b.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: PROPOSED SOLUTION & TECH STACK
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, "SOLUSI TEKNOLOGI MODERN", "Solusi Platform ApotekSim & Full-Stack Tech Stack", 3)

    # Left Solution Card
    add_card(slide3, Inches(0.8), Inches(1.55), Inches(4.5), Inches(5.35), bg_color=COLOR_CARD_BG, border_color=COLOR_TEAL)
    tb_sol = slide3.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(4.1), Inches(4.95))
    tf_s = tb_sol.text_frame
    tf_s.word_wrap = True

    ps1 = tf_s.paragraphs[0]
    ps1.text = "💡 Konsep Solusi ApotekSim"
    ps1.font.size = Pt(17)
    ps1.font.bold = True
    ps1.font.color.rgb = COLOR_CYAN

    ps2 = tf_s.add_paragraph()
    ps2.text = "ApotekSim dirancang sebagai platform web manajemen apotek Cloud-Native yang responsif, mengutamakan validitas data transaksional & rotasi obat terotomatisasi."
    ps2.font.size = Pt(11.5)
    ps2.font.color.rgb = COLOR_WHITE
    ps2.space_before = Pt(8)

    bullets = [
        "🔄 Antrean FEFO Otomatis: Pengeluaran stok mendahulukan batch ED terdekat.",
        "⚡ Realtime WebSockets: Synchronous update stok ke seluruh layar kasir.",
        "🔒 Row Level Security (RLS): Database locking hak akses peran (Admin/Kasir).",
        "🧾 Digital POS & Skrining Resep: Form resep & kalkulasi PPN instan (<2 detik)."
    ]
    for b in bullets:
        pb = tf_s.add_paragraph()
        pb.text = b
        pb.font.size = Pt(10.5)
        pb.font.color.rgb = COLOR_MUTED
        pb.space_before = Pt(10)

    # Right Tech Stack Cards (2x2 Grid)
    techs = [
        {"name": "Next.js 14 & React", "role": "Frontend Framework", "desc": "Server-Side Rendering (SSR) & App Router untuk performa render cepat & SEO responsif.", "badge": "SSR / SSG"},
        {"name": "TypeScript & Tailwind", "role": "Language & Styling", "desc": "Static type-safety mencegah runtime error & styling modern berbasis Shadcn/ui.", "badge": "Type-Safe"},
        {"name": "Supabase & PostgreSQL", "role": "Cloud BaaS & Database", "desc": "PostgreSQL 15, ACID Compliance, Auth JWT, dan WebSockets Realtime Sync.", "badge": "Cloud BaaS"},
        {"name": "Prisma / Drizzle ORM", "role": "Database Mapping", "desc": "Type-safe database query builder & migrasi terstruktur tanpa SQL Injection.", "badge": "Type-Safe DB"}
    ]

    t_w = Inches(3.35)
    t_h = Inches(2.5)

    for i, t in enumerate(techs):
        row = i // 2
        col = i % 2
        tx = Inches(5.6) + col * (t_w + Inches(0.333))
        ty = Inches(1.55) + row * (t_h + Inches(0.35))

        add_card(slide3, tx, ty, t_w, t_h, bg_color=COLOR_ACCENT_BG, border_color=COLOR_CARD_BORDER)

        tb_t = slide3.shapes.add_textbox(tx + Inches(0.2), ty + Inches(0.2), t_w - Inches(0.4), t_h - Inches(0.4))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True

        pt1 = tf_t.paragraphs[0]
        pt1.text = t["name"]
        pt1.font.size = Pt(14)
        pt1.font.bold = True
        pt1.font.color.rgb = COLOR_CYAN

        pt2 = tf_t.add_paragraph()
        pt2.text = f"{t['role']} • [{t['badge']}]"
        pt2.font.size = Pt(9.5)
        pt2.font.bold = True
        pt2.font.color.rgb = COLOR_TEAL
        pt2.space_before = Pt(3)

        pt3 = tf_t.add_paragraph()
        pt3.text = t["desc"]
        pt3.font.size = Pt(10.5)
        pt3.font.color.rgb = COLOR_WHITE
        pt3.space_before = Pt(8)

    # =========================================================================
    # SLIDE 4: DATABASE ARCHITECTURE & SAFETY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, "ARSITEKTUR DATA & KEAMANAN", "Skema Data Relasional & Fitur Keamanan Database", 4)

    # Left Side: Database Tables Schema (4 Card Rows)
    tables = [
        {"title": "Tabel Supplier (Master)", "fields": "ID_Supplier (PK), Nama_Supplier, No_Telepon, Alamat", "rel": "Relasi 1:N ke Pembelian_Obat"},
        {"title": "Tabel Obat (Master Profil)", "fields": "ID_Obat (PK), Nama_Obat, Kategori_Obat, Harga_Jual, Batas_Min, Stok_Total", "rel": "Relasi 1:N ke Detail_Pembelian_Obat"},
        {"title": "Tabel Pembelian_Obat (Restock)", "fields": "ID_Pembelian (PK), No_Faktur, ID_Supplier (FK), Tanggal_Pembelian, Total_Pengeluaran", "rel": "Dokumen Pengadaan Masuk"},
        {"title": "Tabel Detail_Pembelian_Obat (Batch)", "fields": "ID_Detail (PK), ID_Pembelian (FK), ID_Obat (FK), Nomor_Batch, ED, Stok_Sisa", "rel": "Acuan Utama Algoritma FEFO"}
    ]

    tbl_w = Inches(5.65)
    tbl_h = Inches(1.22)
    for i, tbl in enumerate(tables):
        ty = Inches(1.55) + i * (tbl_h + Inches(0.15))
        add_card(slide4, Inches(0.8), ty, tbl_w, tbl_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)

        tb = slide4.shapes.add_textbox(Inches(1.0), ty + Inches(0.12), tbl_w - Inches(0.4), tbl_h - Inches(0.24))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"📁 {tbl['title']}"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CYAN

        p2 = tf.add_paragraph()
        p2.text = f"Kolom: {tbl['fields']}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_WHITE

        p3 = tf.add_paragraph()
        p3.text = f"Fungsi: {tbl['rel']}"
        p3.font.size = Pt(9)
        p3.font.color.rgb = COLOR_MUTED

    # Right Side: Security & FEFO Mechanics Cards
    add_card(slide4, Inches(6.88), Inches(1.55), Inches(5.65), Inches(2.6), bg_color=COLOR_ACCENT_BG, border_color=COLOR_TEAL)
    tb_fefo = slide4.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(2.3))
    tf_f = tb_fefo.text_frame
    tf_f.word_wrap = True

    pf1 = tf_f.paragraphs[0]
    pf1.text = "🔄 Algoritma Antrean FEFO (First Expired, First Out)"
    pf1.font.size = Pt(14.5)
    pf1.font.bold = True
    pf1.font.color.rgb = COLOR_CYAN

    fefo_desc = [
        "1. Pengurutan Batch: Batch diurutkan berdasarkan E(b_1) <= E(b_2) <= ... <= E(b_n).",
        "2. Alokasi Stok Sekuensial: d_i = min(Q_sisa, q_i) memotong batch ED terdekat dahulu.",
        "3. Validasi Atomic: Menolak transaksi jika total stok terkonsolidasi tidak cukup.",
        "4. Mencegah Human Error: Kasir tidak perlu memilih batch manual di rak."
    ]
    for fd in fefo_desc:
        pf = tf_f.add_paragraph()
        pf.text = fd
        pf.font.size = Pt(10)
        pf.font.color.rgb = COLOR_WHITE
        pf.space_before = Pt(4)

    # RLS Card
    add_card(slide4, Inches(6.88), Inches(4.3), Inches(5.65), Inches(2.6), bg_color=COLOR_ACCENT_BG, border_color=COLOR_EMERALD)
    tb_rls = slide4.shapes.add_textbox(Inches(7.1), Inches(4.45), Inches(5.2), Inches(2.3))
    tf_r = tb_rls.text_frame
    tf_r.word_wrap = True

    pr1 = tf_r.paragraphs[0]
    pr1.text = "🛡️ Kebijakan Row Level Security (RLS) PostgreSQL"
    pr1.font.size = Pt(14.5)
    pr1.font.bold = True
    pr1.font.color.rgb = COLOR_EMERALD

    rls_desc = [
        "• Proteksi Level Database: Kebijakan RLS memverifikasi token JWT pengguna.",
        "• Hak Akses Read (SELECT): Semua staf terautentikasi (Kasir/Admin) dapat membaca.",
        "• Hak Akses Mutasi (INSERT/UPDATE/DELETE): Dikunci khusus untuk role Admin.",
        "• Anti-Bypass Client: Kasir yang mencoba manipulasi harga via konsol JS melempar Status Error 403 Forbidden."
    ]
    for rd in rls_desc:
        pr = tf_r.add_paragraph()
        pr.text = rd
        pr.font.size = Pt(10)
        pr.font.color.rgb = COLOR_WHITE
        pr.space_before = Pt(4)

    # =========================================================================
    # SLIDE 5: SYSTEM DIAGRAMS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, "PERANCANGAN STRUKTUR SISTEM", "Diagram Konteks, DFD Level 0, dan ERD Basis Data", 5)

    diagrams = [
        {"file": "diagram_konteks.png", "title": "Diagram Konteks", "desc": "Alur interaksi entitas Admin & Kasir dengan ApotekSim."},
        {"file": "dfd_level0.png", "title": "DFD Level 0", "desc": "Proses pengadaan, stok batch, POS kasir & laporan."},
        {"file": "erd_database.png", "title": "ERD Basis Data", "desc": "Relasi 1:N Supplier, Obat, Pembelian & Batch Detail."}
    ]

    d_w = Inches(3.64)
    d_h = Inches(5.35)
    start_dx = Inches(0.8)
    gap_d = Inches(0.4)

    for i, diag in enumerate(diagrams):
        dx = start_dx + i * (d_w + gap_d)
        dy = Inches(1.55)

        # Container Card
        add_card(slide5, dx, dy, d_w, d_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)

        # Header inside Card
        tb_dh = slide5.shapes.add_textbox(dx + Inches(0.15), dy + Inches(0.15), d_w - Inches(0.3), Inches(0.45))
        tf_dh = tb_dh.text_frame
        p_dh = tf_dh.paragraphs[0]
        p_dh.text = diag["title"]
        p_dh.font.size = Pt(13.5)
        p_dh.font.bold = True
        p_dh.font.color.rgb = COLOR_CYAN
        p_dh.alignment = PP_ALIGN.CENTER

        # Image if exists
        img_path = diag["file"]
        if os.path.exists(img_path):
            max_iw_in = 3.24
            max_ih_in = 3.6
            fw, fh = get_fitted_dimensions(img_path, max_iw_in, max_ih_in)
            
            # Center inside image frame box
            ix = dx + Inches(0.2) + (Inches(max_iw_in) - fw) / 2
            iy = dy + Inches(0.65) + (Inches(max_ih_in) - fh) / 2
            
            add_card(slide5, dx + Inches(0.2), dy + Inches(0.65), Inches(max_iw_in), Inches(max_ih_in), bg_color=RGBColor(255, 255, 255), border_color=None)
            slide5.shapes.add_picture(img_path, ix, iy, fw, fh)

        # Description text
        tb_dd = slide5.shapes.add_textbox(dx + Inches(0.2), dy + Inches(4.35), d_w - Inches(0.4), Inches(0.85))
        tf_dd = tb_dd.text_frame
        tf_dd.word_wrap = True
        p_dd = tf_dd.paragraphs[0]
        p_dd.text = diag["desc"]
        p_dd.font.size = Pt(10.5)
        p_dd.font.color.rgb = COLOR_WHITE
        p_dd.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 6: PRODUCT INTERFACE DEMO (KASIR & KATALOG)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "ANTARMUKA APLIKASI (DEMO PRODUK - BAGIAN 1)", "Demonstrasi UI: Kasir POS & Katalog Obat Reaktif", 6)

    ui_group1 = [
        {
            "file": "ui_pos_kasir.png",
            "title": "🛒 Antarmuka Kasir Digital (Point of Sales)",
            "bullets": [
                "• Instant Search: Pencarian obat cepat & filter kategori.",
                "• Keranjang Reaktif: Perhitungan subtotal, PPN 11% & kembalian.",
                "• Skrining Resep: Form data Dokter & SIP untuk obat keras."
            ]
        },
        {
            "file": "ui_katalog_obat.png",
            "title": "💊 Katalog Obat & Status FEFO Real-Time",
            "bullets": [
                "• Visual Cards: Menampilkan gambar, kategori & stok obat.",
                "• Indicator FEFO: Menampilkan sisa batch aktif & ED terdekat.",
                "• Realtime Sync: Auto-update stok saat terjadi checkout."
            ]
        }
    ]

    ui_w = Inches(5.65)
    ui_h = Inches(5.35)

    for i, ui in enumerate(ui_group1):
        ux = Inches(0.8) + i * (ui_w + Inches(0.433))
        uy = Inches(1.55)

        add_card(slide6, ux, uy, ui_w, ui_h, bg_color=COLOR_CARD_BG, border_color=COLOR_TEAL)

        # Title
        tb_u = slide6.shapes.add_textbox(ux + Inches(0.2), uy + Inches(0.15), ui_w - Inches(0.4), Inches(0.45))
        tf_u = tb_u.text_frame
        p_u = tf_u.paragraphs[0]
        p_u.text = ui["title"]
        p_u.font.size = Pt(13.5)
        p_u.font.bold = True
        p_u.font.color.rgb = COLOR_CYAN

        # Image
        if os.path.exists(ui["file"]):
            max_iw = 5.25
            max_ih = 3.2
            fw, fh = get_fitted_dimensions(ui["file"], max_iw, max_ih)
            ix = ux + Inches(0.2) + (Inches(max_iw) - fw) / 2
            iy = uy + Inches(0.65) + (Inches(max_ih) - fh) / 2

            add_card(slide6, ux + Inches(0.2), uy + Inches(0.65), Inches(max_iw), Inches(max_ih), bg_color=RGBColor(15, 23, 42), border_color=COLOR_CARD_BORDER)
            slide6.shapes.add_picture(ui["file"], ix, iy, fw, fh)

        # Description Bullets
        tb_b = slide6.shapes.add_textbox(ux + Inches(0.2), uy + Inches(3.95), ui_w - Inches(0.4), Inches(1.2))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        for b in ui["bullets"]:
            pb = tf_b.add_paragraph()
            pb.text = b
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = COLOR_WHITE
            pb.space_before = Pt(3)

    # =========================================================================
    # SLIDE 7: PRODUCT INTERFACE DEMO (ADMIN & RESTOCK)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "ANTARMUKA APLIKASI (DEMO PRODUK - BAGIAN 2)", "Demonstrasi UI: Dasbor Admin & Pengadaan Restock", 7)

    ui_group2 = [
        {
            "file": "ui_dashboard_admin.png",
            "title": "📊 Dasbor Executive & Alert Low-Stock",
            "bullets": [
                "• Ringkasan Omset: Grafik penjualan harian & statistik total.",
                "• Alert Berkedip Merah: Peringatan otomatis obat stok minimum.",
                "• Quick Actions: Akses cepat ke tambah obat & supplier."
            ]
        },
        {
            "file": "ui_tabel_restock.png",
            "title": "📦 Manajemen Restock & Batch Expired",
            "bullets": [
                "• Form Pengadaan: Pencatatan faktur masuk dari Supplier.",
                "• Multi-Batch Input: Input Nomor Batch, ED & Harga Beli.",
                "• Auto Consolidation: Memperbarui total stok obat otomatis."
            ]
        }
    ]

    for i, ui in enumerate(ui_group2):
        ux = Inches(0.8) + i * (ui_w + Inches(0.433))
        uy = Inches(1.55)

        add_card(slide7, ux, uy, ui_w, ui_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CYAN)

        # Title
        tb_u = slide7.shapes.add_textbox(ux + Inches(0.2), uy + Inches(0.15), ui_w - Inches(0.4), Inches(0.45))
        tf_u = tb_u.text_frame
        p_u = tf_u.paragraphs[0]
        p_u.text = ui["title"]
        p_u.font.size = Pt(13.5)
        p_u.font.bold = True
        p_u.font.color.rgb = COLOR_CYAN

        # Image
        if os.path.exists(ui["file"]):
            max_iw = 5.25
            max_ih = 3.2
            fw, fh = get_fitted_dimensions(ui["file"], max_iw, max_ih)
            ix = ux + Inches(0.2) + (Inches(max_iw) - fw) / 2
            iy = uy + Inches(0.65) + (Inches(max_ih) - fh) / 2

            add_card(slide7, ux + Inches(0.2), uy + Inches(0.65), Inches(max_iw), Inches(max_ih), bg_color=RGBColor(15, 23, 42), border_color=COLOR_CARD_BORDER)
            slide7.shapes.add_picture(ui["file"], ix, iy, fw, fh)

        # Description Bullets
        tb_b = slide7.shapes.add_textbox(ux + Inches(0.2), uy + Inches(3.95), ui_w - Inches(0.4), Inches(1.2))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        for b in ui["bullets"]:
            pb = tf_b.add_paragraph()
            pb.text = b
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = COLOR_WHITE
            pb.space_before = Pt(3)

    # =========================================================================
    # SLIDE 8: TESTING & CONCLUSION
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, "PENGUJIAN & PENUTUP", "Matriks Pengujian UAT (Black-Box) & Kesimpulan Utama", 8)

    # Left Side: Black-Box Test Matrix Card
    add_card(slide8, Inches(0.8), Inches(1.55), Inches(6.5), Inches(5.35), bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)

    tb_tm = slide8.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(6.1), Inches(4.95))
    tf_tm = tb_tm.text_frame
    tf_tm.word_wrap = True

    ptm1 = tf_tm.paragraphs[0]
    ptm1.text = "🧪 Matriks Hasil Pengujian (Black-Box UAT)"
    ptm1.font.size = Pt(16)
    ptm1.font.bold = True
    ptm1.font.color.rgb = COLOR_CYAN

    tests = [
        {"id": "TC-01", "name": "Autentikasi & RBAC", "res": "Kasir dilarang ke Admin", "status": "PASSED 🟢"},
        {"id": "TC-02", "name": "Validasi Master Obat", "res": "Harga <= 0 ditolak sistem", "status": "PASSED 🟢"},
        {"id": "TC-03", "name": "Transaksi Kasir POS", "res": "Struk & stok terpotong <2s", "status": "PASSED 🟢"},
        {"id": "TC-04", "name": "Skrining Resep Dokter", "res": "Data Dokter/SIP wajib diisi", "status": "PASSED 🟢"},
        {"id": "TC-05", "name": "Otomatisasi Antrean FEFO", "res": "Batch ED terdekat terpotong", "status": "PASSED 🟢"},
        {"id": "TC-06", "name": "Notifikasi Low Stock", "res": "Indikator merah saat stok min", "status": "PASSED 🟢"}
    ]

    for t in tests:
        pt = tf_tm.add_paragraph()
        pt.text = f"• {t['id']} [{t['name']}]: {t['res']} → {t['status']}"
        pt.font.size = Pt(10.5)
        pt.font.color.rgb = COLOR_WHITE
        pt.space_before = Pt(6)

    # Right Side: Key Conclusion Card
    add_card(slide8, Inches(7.533), Inches(1.55), Inches(5.0), Inches(5.35), bg_color=COLOR_ACCENT_BG, border_color=COLOR_EMERALD)

    tb_conc = slide8.shapes.add_textbox(Inches(7.75), Inches(1.75), Inches(4.55), Inches(4.95))
    tf_c = tb_conc.text_frame
    tf_c.word_wrap = True

    pc1 = tf_c.paragraphs[0]
    pc1.text = "🏁 Kesimpulan Utama Proyek"
    pc1.font.size = Pt(16)
    pc1.font.bold = True
    pc1.font.color.rgb = COLOR_EMERALD

    conclusions = [
        "🚀 Efisiensi Transaksi Kasir: Kecepatan transaksi POS <2 detik menekan antrean fisik konsumen.",
        "📦 Eliminasi Kerugian Expired: Metode FEFO terotomatisasi penuh menjamin rotasi obat yang aman.",
        "🔒 Keamanan Cloud & RLS: PostgreSQL RLS di Supabase terbukti andal mengunci akses dari manipulasi client.",
        "📊 Data Real-Time: Integrasi WebSockets memastikan sinkronisasi stok instan antar browser kasir."
    ]

    for c in conclusions:
        pc = tf_c.add_paragraph()
        pc.text = c
        pc.font.size = Pt(10.5)
        pc.font.color.rgb = COLOR_WHITE
        pc.space_before = Pt(10)

    # Save presentation
    output_filename = "presentasi-proyek.pptx"
    prs.save(output_filename)
    print(f"Sukses men-generate {output_filename} dengan desain kreatif, presisi visual & proporsi proporsional!")

if __name__ == "__main__":
    create_presentation()
